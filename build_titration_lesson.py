"""
TLAG Lesson Builder - Proper Template Integration
Injects content into template placeholders, keeps headers intact
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TEMPLATE = "WSO Learn Like A GEM Template (1).pptx"
OUTPUT = "Lesson_4_Titration_Calculations.pptx"

# Colors
BLACK = RGBColor(0, 0, 0)
LIGHT_BLUE = RGBColor(0, 112, 192)
RED = RGBColor(192, 0, 0)
GREEN = RGBColor(0, 128, 0)

# Layout indices
LAYOUTS = {
    "DO_NOW": 1,
    "OUTCOMES": 2,
    "I_DO": 4,
    "WE_DO": 5,
    "YOU_DO": 6,
    "CHECK": 7,
    "EXIT": 8
}

def set_placeholder_text(slide, idx, content):
    """Set text in a placeholder. Content: list of (text, bold, color) tuples."""
    try:
        ph = slide.placeholders[idx]
        # Reposition the placeholder higher on the slide
        ph.top = Inches(0.85)  # Move up closer to header
        ph.left = Inches(0.3)
        ph.width = Inches(12.7)
        ph.height = Inches(6.3)
    except KeyError:
        return
    
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    
    for item in content:
        text, bold, color = item
        if text.startswith("\n"):
            p = tf.add_paragraph()
            text = text[1:]
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(24)  # Reduced from 27 to fit content
        run.font.bold = bold
        run.font.color.rgb = color if color else BLACK

def remove_title(slide):
    """Remove the title placeholder (idx=0) from the slide."""
    try:
        title_ph = slide.placeholders[0]
        sp = title_ph._element
        sp.getparent().remove(sp)
    except (KeyError, Exception):
        pass

def build():
    print("🚀 Building TLAG Lesson (Template Placeholders)...")
    prs = Presentation(TEMPLATE)
    
    # Clear existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # === SLIDE 1: Do Now ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["DO_NOW"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Complete the retrieval questions in silence. (6 mins)", True, None),
        ("\n", False, None),
        ("\n1. How many cm³ are in 1 dm³? ", False, None), ("1000", True, RED),
        ("\n2. Convert 25.0 cm³ into dm³. ", False, None), ("0.025 dm³", True, RED),
        ("\n3. Convert 500 cm³ into dm³. ", False, None), ("0.500 dm³", True, RED),
        ("\n4. Formula: Moles, Concentration, Volume. ", False, None), ("n = C × V", True, RED),
        ("\n5. Units for Concentration? ", False, None), ("mol/dm³", True, RED),
        ("\n6. Moles in 0.5 dm³ of 2.0 mol/dm³ HCl? ", False, None), ("1.0 mol", True, RED),
        ("\n7. Balance: NaOH + H₂SO₄ → ", False, None), ("2NaOH + H₂SO₄ → Na₂SO₄ + 2H₂O", True, RED),
        ("\n8. Moles NaOH per 1 mol H₂SO₄? ", False, None), ("2", True, RED),
        ("\n9. (Stretch) Mr of NaOH? ", False, LIGHT_BLUE), ("40", True, RED),
        ("\n10. (Stretch) Mass of 0.5 mol NaOH? ", False, LIGHT_BLUE), ("20g", True, RED),
    ])

    # === SLIDE 2: Learning Outcomes ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["OUTCOMES"]])
    remove_title(s)
    # Left body (idx=1) - Learning Outcome
    set_placeholder_text(s, 1, [
        ("Calculate the unknown concentration of a solution using titration results and balanced symbol equations.", False, None),
    ])
    # Right body (idx=2) - To Know
    set_placeholder_text(s, 2, [
        ("1. ", False, None), ("Concentration", True, LIGHT_BLUE), (" is moles per dm³.", False, None),
        ("\n2. ", False, None), ("Stoichiometry", True, LIGHT_BLUE), (" is the ratio of moles.", False, None),
        ("\n3. ", False, None), ("Titration", True, LIGHT_BLUE), (" finds unknown concentration.", False, None),
        ("\n4. ", False, None), ("Golden Rule:", True, RED), (" cm³ ÷ 1000 = dm³", False, None),
        ("\n5. ", False, None), ("Extension:", True, LIGHT_BLUE), (" mol/dm³ × Mr = g/dm³", False, None),
    ])

    # === SLIDE 3: I Do - Grid Method ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["I_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("The R-V-C-n Grid Method", True, None),
        ("\n", False, None),
        ("\nProblem:", True, RED), (" Titration questions hide numbers in paragraphs.", False, None),
        ("\nSolution:", True, GREEN), (" Use the R-V-C-n Grid.", False, None),
        ("\n", False, None),
        ("\nStrategy:", True, None),
        ("\n1. ", False, None), ("Ratio:", True, LIGHT_BLUE), (" From equation.", False, None),
        ("\n2. ", False, None), ("Volume:", True, LIGHT_BLUE), (" Convert ÷ 1000.", False, None),
        ("\n3. ", False, None), ("Moles:", True, LIGHT_BLUE), (" n = C × V.", False, None),
        ("\n4. ", False, None), ("Bridge:", True, RED), (" Use ratio.", False, None),
        ("\n5. ", False, None), ("Conc:", True, LIGHT_BLUE), (" C = n ÷ V.", False, None),
    ])

    # === SLIDE 4: I Do - 1:1 Ratio ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["I_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("1:1 Ratio Example", True, None),
        ("\nProblem:", True, None), (" 25.0 cm³ of 0.1 mol/dm³ NaOH + 20.0 cm³ HCl. Find [HCl].", False, None),
        ("\nEquation:", True, None), (" HCl + NaOH → NaCl + H₂O", True, LIGHT_BLUE),
        ("\n", False, None),
        ("\nStep 1:", True, None), (" n = 0.1 × 0.025 = ", False, None), ("0.0025 mol", True, GREEN),
        ("\nStep 2 (1:1):", True, None), (" HCl = ", False, None), ("0.0025 mol", True, GREEN),
        ("\nStep 3:", True, None), (" C = 0.0025 ÷ 0.020 = ", False, None), ("0.125 mol/dm³", True, GREEN),
    ])

    # === SLIDE 5: I Do - 1:2 Ratio ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["I_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("1:2 Ratio – The Grade 9 Trap!", True, RED),
        ("\nProblem:", True, None), (" 25.0 cm³ NaOH + 20.0 cm³ of 0.5 mol/dm³ H₂SO₄. Find [NaOH].", False, None),
        ("\nEquation:", True, None), (" H₂SO₄ + ", False, LIGHT_BLUE), ("2", True, RED), ("NaOH → Na₂SO₄ + 2H₂O", False, LIGHT_BLUE),
        ("\n", False, None),
        ("\nStep 1:", True, None), (" n = 0.5 × 0.020 = ", False, None), ("0.01 mol", True, GREEN),
        ("\nStep 2 ", True, None), ("(×2!):", True, RED), (" NaOH = 0.01 × 2 = ", False, None), ("0.02 mol", True, GREEN),
        ("\nStep 3:", True, None), (" C = 0.02 ÷ 0.025 = ", False, None), ("0.8 mol/dm³", True, GREEN),
    ])

    # === SLIDE 6: We Do ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["WE_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Scaffolded Practice - Draw the Grid!", True, None),
        ("\n", False, None),
        ("\nQuestion:", True, None), (" 24.0 cm³ of 0.2 mol/dm³ KOH + 30.0 cm³ HNO₃. Find [HNO₃].", False, None),
        ("\nEquation:", True, None), (" KOH + HNO₃ → KNO₃ + H₂O", True, LIGHT_BLUE),
        ("\n", False, None),
        ("\nAnswer:", True, None),
        ("\n• Ratio: 1:1", False, None),
        ("\n• n(KOH) = 0.2 × 0.024 = ", False, None), ("0.0048 mol", True, GREEN),
        ("\n• n(HNO₃) = 0.0048 mol", False, None),
        ("\n• C = 0.0048 ÷ 0.030 = ", False, None), ("0.16 mol/dm³", True, GREEN),
    ])

    # === SLIDE 7: I Do Extension ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["I_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Converting to g/dm³", True, None),
        ("\n", False, None),
        ("\nLogic:", True, LIGHT_BLUE), (" Moles count particles. Grams measure mass.", False, None),
        ("\nFormula:", True, LIGHT_BLUE), (" Conc (g/dm³) = Conc (mol/dm³) × Mr", False, None),
        ("\n", False, None),
        ("\nExample:", True, None),
        ("\n• Concentration = 0.5 mol/dm³", False, None),
        ("\n• Substance = NaOH (Mr = 40)", False, None),
        ("\n• 0.5 × 40 = ", False, None), ("20 g/dm³", True, GREEN),
    ])

    # === SLIDE 8: You Do ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["YOU_DO"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Independent Practice", True, None),
        ("\n", False, None),
        ("\nTask 1 (Grade 5):", True, None), (" 25 cm³ 1.0 mol/dm³ NaOH + 50 cm³ HCl. [HCl]?", False, None),
        ("\nAnswer: ", False, None), ("0.5 mol/dm³", True, GREEN),
        ("\n", False, None),
        ("\nTask 2 (Grade 7):", True, None), (" 20 cm³ 0.1 mol/dm³ H₂SO₄ + 25 cm³ NaOH. [NaOH]?", False, None),
        ("\nAnswer: ", False, None), ("0.16 mol/dm³", True, GREEN),
        ("\n", False, None),
        ("\nTask 3 (Grade 9):", True, None), (" Convert Task 2 to g/dm³ (Mr=40).", False, None),
        ("\nAnswer: ", False, None), ("6.4 g/dm³", True, GREEN),
    ])

    # === SLIDE 9: Affirmative Checking ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["CHECK"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Exam Question (5 Marks)", True, None),
        ("\n25 cm³ NH₃ + 30 cm³ 0.2 mol/dm³ H₂SO₄. [NH₃] in g/dm³? (N=14, H=1)", False, None),
        ("\n2NH₃ + H₂SO₄ → (NH₄)₂SO₄", True, LIGHT_BLUE),
        ("\n", False, None),
        ("\n1. n(Acid) = 0.2 × 0.030 = ", False, None), ("0.006 mol ✓", True, GREEN),
        ("\n2. Ratio 1:2 → n(NH₃) = ", False, None), ("0.012 mol ✓", True, GREEN),
        ("\n3. C = 0.012 ÷ 0.025 = ", False, None), ("0.48 mol/dm³ ✓", True, GREEN),
        ("\n4. Mr = 14 + 3 = ", False, None), ("17 ✓", True, GREEN),
        ("\n5. 0.48 × 17 = ", False, None), ("8.16 g/dm³ ✓", True, GREEN),
    ])

    # === SLIDE 10: Exit Ticket ===
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["EXIT"]])
    remove_title(s)
    set_placeholder_text(s, 1, [
        ("Answer on your post-it:", True, None),
        ("\n", False, None),
        ("\nA student uses '25' instead of '0.025' (forgot to convert).", False, None),
        ("\nHow is their answer affected?", False, None),
        ("\n", False, None),
        ("\nA) 1000× too big", False, None),
        ("\nB) 1000× too small", False, None),
        ("\nC) Correct", False, None),
        ("\n", False, None),
        ("\nAnswer: ", True, None), ("B) 1000× too small", True, GREEN),
    ])

    prs.save(OUTPUT)
    print(f"✅ Saved: {OUTPUT}")
    print("   10 slides with template headers intact")

if __name__ == "__main__":
    build()
