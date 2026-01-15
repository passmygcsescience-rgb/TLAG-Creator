import streamlit as st
import json
import os
from openai import OpenAI
from pptx import Presentation

# ==============================================================================
# ⚠️  USER CONFIGURATION REQUIRED  ⚠️
# Run 'audit_template.py' first. Look at the output.
# Update the integers below to match the Layout Indices from YOUR template.
# ==============================================================================
LAYOUT_MAP = {
    "TITLE": 0,          # Title
    "DO_NOW": 1,         # Do Now
    "OUTCOMES": 2,       # Learning Outcomes
    "I_DO": 3,           # I Do
    "WE_DO": 4,          # We Do
    "YOU_DO": 5,         # You Do
    "CHECK": 7,          # Affirmative Checking
    "EXIT": 8            # Exit Tickets
}
# ==============================================================================

TEMPLATE_FILE = "WSO Learn Like A GEM Template (1).pptx"

def get_ai_content(client, topic, spec):
    """Generates curriculum-aligned content in strict JSON format."""
    system_prompt = f"""
    You are an expert Science Teacher for {spec}. Plan a lesson on '{topic}' following the 'Teach Like a GEM' (TLAG) pedagogy.
    
    Output ONLY valid JSON (no markdown) with the following keys:
    {{
        "title_slide": {{ "title": "Lesson Title", "subtitle": "Specification Reference" }},
        "do_now": {{ "title": "Do Now", "body": "3 Retrieval Questions..." }},
        "outcomes": {{ "title": "Learning Outcomes", "body": "To Know: ... \\nTo Be Able To: ..." }},
        "i_do": {{ "title": "I Do (Teacher Model)", "body": "Clear explanation text..." }},
        "we_do": {{ "title": "We Do (Guided)", "body": "Worked example..." }},
        "you_do": {{ "title": "You Do (Independent)", "body": "Bronze/Silver/Gold questions..." }},
        "check": {{ "title": "Affirmative Checking", "body": "Hinge Question..." }},
        "exit": {{ "title": "Exit Ticket", "body": "Final diagnostic question..." }}
    }}
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[{"role": "system", "content": system_prompt}],
            response_format={"type": "json_object"}
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        st.error(f"AI Generation Error: {e}")
        return None

def inject_slide(prs, layout_key, content):
    """Creates a slide and injects text strictly into existing placeholders."""
    layout_idx = LAYOUT_MAP.get(layout_key, 0)
    
    # Safety check
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = 0
        
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # 1. Inject Title
    if slide.shapes.title:
        slide.shapes.title.text = content.get("title", "")
        
    # 2. Inject Body (Find first text placeholder that isn't the title)
    for ph in slide.placeholders:
        if ph.has_text_frame and ph != slide.shapes.title:
            ph.text_frame.text = content.get("body", "")
            break

def build_lesson(api_key, topic, spec):
    if not os.path.exists(TEMPLATE_FILE):
        st.error("❌ Template file not found. Please drop it in this folder.")
        return None

    client = OpenAI(api_key=api_key)
    data = get_ai_content(client, topic, spec)
    
    if not data:
        return None

    prs = Presentation(TEMPLATE_FILE)
    
    # Build the 8-Stage TLAG Sequence
    inject_slide(prs, "TITLE", data["title_slide"])
    inject_slide(prs, "DO_NOW", data["do_now"])
    inject_slide(prs, "OUTCOMES", data["outcomes"])
    inject_slide(prs, "I_DO", data["i_do"])
    inject_slide(prs, "WE_DO", data["we_do"])
    inject_slide(prs, "YOU_DO", data["you_do"])
    inject_slide(prs, "CHECK", data["check"])
    inject_slide(prs, "EXIT", data["exit"])
    
    output_file = f"Lesson_{topic.replace(' ', '_')}.pptx"
    prs.save(output_file)
    return output_file

# --- STREAMLIT UI ---
st.set_page_config(page_title="GEMS TLAG Builder", layout="centered")
st.title("💎 GEMS TLAG Lesson Builder")
st.caption("Automated Science Lesson Generator (AQA/IB)")

if not os.path.exists(TEMPLATE_FILE):
    st.warning(f"⚠️ Waiting for `{TEMPLATE_FILE}`")
    st.info("Drag and drop the PowerPoint template into this folder to proceed.")
    st.stop()

with st.sidebar:
    st.header("Settings")
    api_key = st.text_input("OpenAI API Key", type="password")
    spec = st.selectbox("Specification", ["AQA GCSE Chemistry", "AQA GCSE Physics", "AQA GCSE Biology", "IB Chemistry HL", "IB Chemistry SL"])

topic = st.text_input("Lesson Topic", placeholder="e.g. Fractional Distillation")

if st.button("Generate Lesson", type="primary"):
    if not api_key or not topic:
        st.error("Please enter both an API Key and a Topic.")
    else:
        with st.spinner("💎 Mining pedagogical crystals..."):
            result = build_lesson(api_key, topic, spec)
            if result:
                st.success("Lesson Generated Successfully!")
                with open(result, "rb") as f:
                    st.download_button("📥 Download PowerPoint", f, file_name=result)
