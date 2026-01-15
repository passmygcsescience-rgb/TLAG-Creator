"""
TLAG PowerPoint Builder
Generates GEMS TLAG lesson PowerPoints with proper layout and images.
Fixed font sizing and image integration.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
from io import BytesIO

class TLAGPowerPointBuilder:
    """Builds TLAG-compliant PowerPoint presentations with images."""
    
    # Template layouts (matching the TLAG template slide master)
    # Layout indices based on order in template:
    # Row 1: Title Slide, Do Now, Learning Outcomes, I Do We Do You Do, I Do
    # Row 2: We Do, You Do, Affirmative Checking, Exit Tickets, Empty
    # Row 3: Empty with Header, Turn Talk, Show Call, Show Me, Everybody Writes
    # Row 4: Cold Call, Right is Right, Stretch It, Habits of Discussion, What to Do
    # Row 5: Whole Class Reset, Teacher Radar, Secure Student attention, Blank
    LAYOUTS = {
        # Core TLAG slides
        "TITLE": 0,
        "DO_NOW": 1,
        "OUTCOMES": 2,
        "I_DO_WE_DO_YOU_DO": 3,
        "I_DO": 4,
        "WE_DO": 5,
        "YOU_DO": 6,
        "CHECK": 7,              # Affirmative Checking
        "EXIT": 8,               # Exit Tickets
        "EMPTY": 9,
        # TLAG Techniques
        "EMPTY_HEADER": 10,
        "TURN_TALK": 11,         # Turn and Talk
        "SHOW_CALL": 12,
        "SHOW_ME": 13,
        "EVERYBODY_WRITES": 14,
        "COLD_CALL": 15,
        "RIGHT_IS_RIGHT": 16,
        "STRETCH_IT": 17,
        "HABITS_DISCUSSION": 18,
        "WHAT_TO_DO": 19,
        "WHOLE_CLASS_RESET": 20,
        "TEACHER_RADAR": 21,
        "SECURE_ATTENTION": 22,
        "BLANK": 23
    }
    
    # Colors
    BLACK = RGBColor(0, 0, 0)
    RED = RGBColor(192, 0, 0)
    GREEN = RGBColor(0, 128, 0)
    BLUE = RGBColor(0, 112, 192)
    PURPLE = RGBColor(102, 45, 145)
    ORANGE = RGBColor(255, 102, 0)
    GRAY = RGBColor(100, 100, 100)
    
    # Slide dimensions (standard 16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Image mapping - maps topic keywords to image filenames
    IMAGE_MAPPING = {
        # Atomic structure
        "atom": "atomic_structure/atom_structure.png",
        "atoms": "atomic_structure/atom_structure.png",
        "element": "atomic_structure/atom_structure.png",
        "compound": "atomic_structure/atom_structure.png",
        "periodic table": "atomic_structure/periodic_table.png",
        "group": "atomic_structure/periodic_table.png",
        "period": "atomic_structure/periodic_table.png",
        "electron": "atomic_structure/electronic_structure.png",
        "electronic": "atomic_structure/electronic_structure.png",
        "shell": "atomic_structure/electronic_structure.png",
        "proton": "atomic_structure/subatomic_particles.png",
        "neutron": "atomic_structure/subatomic_particles.png",
        "nucleus": "atomic_structure/subatomic_particles.png",
        "subatomic": "atomic_structure/subatomic_particles.png",
        "alkali metal": "atomic_structure/alkali_metals.png",
        "lithium": "atomic_structure/alkali_metals.png",
        "sodium": "atomic_structure/alkali_metals.png",
        "potassium": "atomic_structure/alkali_metals.png",
        "halogen": "atomic_structure/halogens.png",
        "chlorine": "atomic_structure/halogens.png",
        "bromine": "atomic_structure/halogens.png",
        "noble gas": "atomic_structure/noble_gases.png",
        "helium": "atomic_structure/noble_gases.png",
        "neon": "atomic_structure/noble_gases.png",
        "argon": "atomic_structure/noble_gases.png",
        "transition metal": "atomic_structure/transition_metals.png",
        "iron": "atomic_structure/transition_metals.png",
        "copper": "atomic_structure/transition_metals.png",
        "mixture": "atomic_structure/mixtures_separation.png",
        "separation": "atomic_structure/mixtures_separation.png",
        "filtration": "atomic_structure/mixtures_separation.png",
        "distillation": "atomic_structure/mixtures_separation.png",
        # Bonding
        "ionic": "bonding/ionic_bonding.png",
        "ion": "bonding/ionic_bonding.png",
        "covalent": "bonding/covalent_bonding.png",
        "molecule": "bonding/covalent_bonding.png",
        "metallic": "bonding/metallic_bonding.png",
        "diamond": "bonding/diamond_graphite.png",
        "graphite": "bonding/diamond_graphite.png",
        "graphene": "bonding/diamond_graphite.png",
        "fullerene": "bonding/fullerenes.png",
        "nanotube": "bonding/fullerenes.png",
        "buckyball": "bonding/fullerenes.png",
        "state": "bonding/states_of_matter.png",
        "solid": "bonding/states_of_matter.png",
        "liquid": "bonding/states_of_matter.png",
        "gas": "bonding/states_of_matter.png",
        "particle": "bonding/states_of_matter.png",
        # Quantitative
        "mole": "quantitative/moles_calculation.png",
        "mass": "quantitative/moles_calculation.png",
        "calculation": "quantitative/moles_calculation.png",
        "titration": "quantitative/titration_setup.png",
        "burette": "quantitative/titration_setup.png",
        # Chemical changes
        "reactivity": "chemical_changes/reactivity_series.png",
        "displacement": "chemical_changes/reactivity_series.png",
        "electrolysis": "chemical_changes/electrolysis.png",
        "electrode": "chemical_changes/electrolysis.png",
        "acid": "chemical_changes/acid_base_reaction.png",
        "base": "chemical_changes/acid_base_reaction.png",
        "neutralisation": "chemical_changes/acid_base_reaction.png",
        "neutralization": "chemical_changes/acid_base_reaction.png",
        "ph": "chemical_changes/acid_base_reaction.png",
        # Energy changes
        "exothermic": "energy_changes/exothermic_endothermic.png",
        "endothermic": "energy_changes/exothermic_endothermic.png",
        "energy": "energy_changes/exothermic_endothermic.png",
        "reaction profile": "energy_changes/exothermic_endothermic.png",
        # Rates
        "collision": "rates/collision_theory.png",
        "rate": "rates/collision_theory.png",
        "catalyst": "rates/catalyst.png",
        "enzyme": "rates/catalyst.png",
        "equilibrium": "rates/equilibrium.png",
        "reversible": "rates/equilibrium.png",
        # Organic
        "crude oil": "organic/fractional_distillation.png",
        "fraction": "organic/fractional_distillation.png",
        "hydrocarbon": "organic/alkanes_alkenes.png",
        "alkane": "organic/alkanes_alkenes.png",
        "alkene": "organic/alkanes_alkenes.png",
        "polymer": "organic/polymers.png",
        "polymerisation": "organic/polymers.png",
        # Analysis
        "chromatography": "analysis/chromatography.png",
        "rf": "analysis/chromatography.png",
        "flame test": "analysis/flame_tests.png",
        "gas test": "analysis/gas_tests.png",
        # Atmosphere
        "atmosphere": "atmosphere/atmosphere_composition.png",
        "greenhouse": "atmosphere/greenhouse_effect.png",
        "climate": "atmosphere/greenhouse_effect.png",
        "carbon dioxide": "atmosphere/greenhouse_effect.png",
        # Resources
        "water treatment": "resources/water_treatment.png",
        "potable": "resources/water_treatment.png",
        "haber": "resources/haber_process.png",
        "ammonia": "resources/haber_process.png",
        "fertiliser": "resources/haber_process.png",
        "rust": "resources/corrosion_rusting.png",
        "corrosion": "resources/corrosion_rusting.png",
    }
    
    def __init__(self, template_path: str):
        """Initialize with template path."""
        self.template_path = template_path
        self.prs = None
        self.current_topic = ""
        self.base_image_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "images", "chemistry")
        
    def create_presentation(self):
        """Create a new presentation from template."""
        self.prs = Presentation(self.template_path)
        # Clear existing slides
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        return self
    
    def _add_slide(self, layout_name: str):
        """Add a slide with the specified layout."""
        layout_idx = self.LAYOUTS.get(layout_name, 1)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        self._remove_all_placeholders(slide)
        return slide
    
    def _remove_all_placeholders(self, slide):
        """Remove ALL placeholders from slide to prevent 'Click to add text' boxes."""
        # Collect all placeholder elements first
        placeholders_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholders_to_remove.append(shape)
        
        # Remove them
        for ph in placeholders_to_remove:
            try:
                sp = ph._element
                sp.getparent().remove(sp)
            except:
                pass
    
    def _find_image_for_text(self, text: str) -> str:
        """Find appropriate image based on text content."""
        text_lower = text.lower()
        for keyword, image_path in self.IMAGE_MAPPING.items():
            if keyword in text_lower:
                full_path = os.path.join(self.base_image_path, image_path)
                if os.path.exists(full_path):
                    return full_path
        return None
    
    def _add_image_to_slide(self, slide, image_path: str, left: float = 9.5, top: float = 1.5, width: float = 3.5):
        """Add an image to the slide at specified position."""
        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    width=Inches(width)
                )
                return True
            except Exception as e:
                print(f"Error adding image: {e}")
        return False
    
    def _create_text_box(self, slide, left: float, top: float, width: float, height: float):
        """Create a text box on the slide."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        return txBox.text_frame
    
    def _add_text_content(self, text_frame, content_items: list, font_size: int = 18):
        """Add formatted text content to a text frame with proper sizing."""
        text_frame.word_wrap = True
        
        # Clear existing
        for para in text_frame.paragraphs:
            para.clear()
        
        p = text_frame.paragraphs[0]
        first = True
        
        for item in content_items:
            text, bold, color = item
            
            # Handle newlines by creating new paragraphs
            if text.startswith("\n"):
                if not first:
                    p = text_frame.add_paragraph()
                text = text.lstrip("\n")
                first = False
            
            if text:
                run = p.add_run()
                run.text = text
                run.font.name = "Arial"
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = color
                else:
                    run.font.color.rgb = self.BLACK
    
    def add_do_now(self, questions: list, answers: list):
        """Add Do Now slide with retrieval questions - FIXED LAYOUT."""
        slide = self._add_slide("DO_NOW")
        
        # Create text box with proper sizing
        tf = self._create_text_box(slide, 0.5, 1.3, 8.5, 5.7)
        
        content = [
            ("Complete in silence (6 minutes)", False, self.GRAY),
            ("\n", False, None)
        ]
        
        for i, (q, a) in enumerate(zip(questions, answers), 1):
            content.append((f"\n{i}. {q}", False, None))
            content.append(("\n    → ", False, None))
            content.append((str(a), True, self.RED))
        
        self._add_text_content(tf, content, font_size=18)
        
        # Try to add a relevant image
        combined_text = " ".join(questions)
        image_path = self._find_image_for_text(combined_text)
        if image_path:
            self._add_image_to_slide(slide, image_path, left=9.5, top=1.5, width=3.3)
        
        return self
    
    def add_outcomes(self, outcome: str, to_know: list):
        """Add Learning Outcomes slide - TWO COLUMN LAYOUT."""
        slide = self._add_slide("OUTCOMES")
        
        # Left column - Learning Outcome
        tf_left = self._create_text_box(slide, 0.5, 1.3, 6.0, 5.7)
        content_left = [
            (f"{outcome}", False, None)
        ]
        self._add_text_content(tf_left, content_left, font_size=20)
        
        # Right column - To Know
        tf_right = self._create_text_box(slide, 6.8, 1.3, 6.0, 5.7)
        content_right = []
        for i, item in enumerate(to_know[:8], 1):  # Limit to 8 items
            # Truncate long items
            display_item = item[:100] + "..." if len(item) > 100 else item
            content_right.append((f"\n{i}. {display_item}", False, None))
        
        self._add_text_content(tf_right, content_right, font_size=16)
        
        return self
    
    def add_i_do_slides(self, i_do_data: dict):
        """Add I Do slides with INFORMATION-BASED content for teacher modelling."""
        title = i_do_data.get("title", "Teacher Modelling")
        content_list = i_do_data.get("content", [])
        examples = i_do_data.get("examples", [])
        key_points = i_do_data.get("key_points", [])
        definitions = i_do_data.get("definitions", {})
        facts = i_do_data.get("facts", [])
        
        # Find image based on title
        image_path = self._find_image_for_text(title + " " + " ".join(content_list[:3]))
        
        # ===== SLIDE 1: I DO - INFORMATION SLIDE =====
        slide1 = self._add_slide("I_DO")
        
        # Determine text width based on whether we have an image
        text_width = 8.0 if image_path else 12.0
        
        tf1 = self._create_text_box(slide1, 0.5, 1.3, text_width, 5.7)
        content1 = [
            (title, True, self.PURPLE),
            ("\n", False, None)
        ]
        
        # Add main content as educational information
        for item in content_list[:6]:  # Limit to 6 items
            # Truncate if too long
            display_item = item[:120] + "..." if len(item) > 120 else item
            content1.append((f"\n• {display_item}", False, None))
        
        self._add_text_content(tf1, content1, font_size=16)
        
        # Add image if found
        if image_path:
            self._add_image_to_slide(slide1, image_path, left=9.0, top=1.2, width=3.8)
        
        # ===== SLIDE 2: KEY DEFINITIONS (if available) =====
        if definitions:
            slide_def = self._add_slide("I_DO")
            tf_def = self._create_text_box(slide_def, 0.5, 1.3, 12.0, 5.7)
            
            content_def = [("Key Definitions", True, self.PURPLE), ("\n", False, None)]
            
            for term, definition in list(definitions.items())[:5]:  # Max 5 definitions
                term_text = str(term)[:40]
                def_text = str(definition)[:150]
                content_def.append((f"\n{term_text}: ", True, self.BLUE))
                content_def.append((def_text, False, None))
            
            self._add_text_content(tf_def, content_def, font_size=16)
        
        # ===== SLIDE 3: KEY FACTS (if available) =====
        if facts:
            slide_facts = self._add_slide("I_DO")
            image_path2 = self._find_image_for_text(" ".join(facts[:2]))
            text_width2 = 8.0 if image_path2 else 12.0
            
            tf_facts = self._create_text_box(slide_facts, 0.5, 1.3, text_width2, 5.7)
            
            content_facts = [("Key Facts to Remember", True, self.PURPLE), ("\n", False, None)]
            
            for fact in facts[:6]:  # Max 6 facts
                fact_text = fact[:120] + "..." if len(fact) > 120 else fact
                content_facts.append((f"\n✓ {fact_text}", False, None))
            
            self._add_text_content(tf_facts, content_facts, font_size=16)
            
            if image_path2:
                self._add_image_to_slide(slide_facts, image_path2, left=9.0, top=1.5, width=3.5)
        
        # ===== SLIDE 4: WORKED EXAMPLE (if available) =====
        if examples:
            slide2 = self._add_slide("I_DO")
            tf2 = self._create_text_box(slide2, 0.5, 1.3, 12.0, 5.7)
            
            content2 = [("Worked Example", True, self.PURPLE), ("\n", False, None)]
            
            for i, example in enumerate(examples[:2], 1):  # Max 2 examples
                problem = example.get('problem', example.get('question', ''))
                content2.append((f"\n{problem}", True, None))
                
                steps = example.get('steps', [])[:5]  # Max 5 steps
                for j, step in enumerate(steps, 1):
                    step_text = step[:100] + "..." if len(step) > 100 else step
                    content2.append((f"\n   Step {j}: {step_text}", False, None))
                
                if example.get('answer'):
                    content2.append(("\n   ✓ Answer: ", False, None))
                    content2.append((str(example.get('answer'))[:80], True, self.GREEN))
                content2.append(("\n", False, None))
            
            self._add_text_content(tf2, content2, font_size=15)
        
        # ===== SLIDE 5: KEY POINTS SUMMARY (if available) =====
        if key_points:
            slide3 = self._add_slide("I_DO")
            tf3 = self._create_text_box(slide3, 0.5, 1.3, 8.0, 5.7)
            
            content3 = [("Remember These Key Points", True, self.PURPLE), ("\n", False, None)]
            
            for point in key_points[:6]:  # Max 6 points
                point_text = point[:100] + "..." if len(point) > 100 else point
                content3.append((f"\n★ {point_text}", False, None))
            
            self._add_text_content(tf3, content3, font_size=17)
            
            # Add image on key points slide too
            if image_path:
                self._add_image_to_slide(slide3, image_path, left=9.0, top=1.5, width=3.5)
        
        return self
    
    def add_i_do(self, title: str, content_items: list):
        """Legacy single I Do slide."""
        self.add_i_do_slides({
            "title": title,
            "content": content_items
        })
        return self
    
    def add_we_do_slides(self, we_do_data: dict):
        """Add We Do slides - Uses template layouts including TURN_TALK."""
        examples = we_do_data.get("examples", [])
        
        # Convert legacy format
        if not examples and we_do_data.get("question"):
            examples = [{
                "question": we_do_data.get("question"),
                "steps": we_do_data.get("steps", []),
                "answer": we_do_data.get("answer", "")
            }]
        
        # Find image
        combined_text = " ".join([e.get("question", "") for e in examples])
        image_path = self._find_image_for_text(combined_text)
        
        # ===== SLIDE 1: WE DO OVERVIEW (uses WE_DO template) =====
        slide1 = self._add_slide("WE_DO")
        text_width = 8.0 if image_path else 12.0
        
        tf1 = self._create_text_box(slide1, 0.5, 1.3, text_width, 5.7)
        content1 = []
        
        if examples:
            ex = examples[0]
            q_text = ex.get('question', '')[:150]
            content1.append((f"📝 {q_text}", True, None))
            
            for i, step in enumerate(ex.get('steps', [])[:5], 1):
                step_text = step[:80] + "..." if len(step) > 80 else step
                content1.append((f"\n   Step {i}: {step_text}", False, None))
            
            if ex.get('answer'):
                content1.append(("\n   ✓ ", False, None))
                content1.append((str(ex.get('answer'))[:60], True, self.GREEN))
        
        if content1:
            self._add_text_content(tf1, content1, font_size=16)
        
        if image_path:
            self._add_image_to_slide(slide1, image_path, left=9.0, top=1.5, width=3.5)
        
        # ===== SLIDE 2: TURN AND TALK (uses TURN_TALK template) =====
        if len(examples) > 1:
            slide_tt = self._add_slide("TURN_TALK")
            tf_tt = self._create_text_box(slide_tt, 0.5, 1.3, 12.0, 5.7)
            
            ex = examples[1]
            content_tt = [
                (f"📝 {ex.get('question', '')[:120]}", True, None)
            ]
            
            for i, step in enumerate(ex.get('steps', [])[:4], 1):
                step_text = step[:70] + "..." if len(step) > 70 else step
                content_tt.append((f"\n   Step {i}: {step_text}", False, None))
            
            if ex.get('answer'):
                content_tt.append(("\n\n   ✓ Answer: ", False, None))
                content_tt.append((str(ex.get('answer'))[:50], True, self.GREEN))
            
            self._add_text_content(tf_tt, content_tt, font_size=16)
        
        # Additional slides for more examples (max 2 more)
        for i, example in enumerate(examples[1:3], 2):
            slide_n = self._add_slide("WE_DO")
            tf_n = self._create_text_box(slide_n, 0.5, 1.3, 12.0, 5.7)
            
            content_n = [
                (f"👥 WE DO: Example {i}", True, self.BLUE),
                ("\n🗣️ Discuss with your partner", True, self.ORANGE),
                ("\n", False, None)
            ]
            
            q_text = example.get('question', '')[:150]
            content_n.append((f"\n📝 {q_text}", True, None))
            
            for j, step in enumerate(example.get('steps', [])[:5], 1):
                step_text = step[:80] + "..." if len(step) > 80 else step
                content_n.append((f"\n   Step {j}: {step_text}", False, None))
            
            if example.get('answer'):
                content_n.append(("\n\n   ✓ Answer: ", False, None))
                content_n.append((str(example.get('answer'))[:60], True, self.GREEN))
            
            self._add_text_content(tf_n, content_n, font_size=16)
        
        return self
    
    def add_we_do(self, title: str, question: str, steps: list):
        """Legacy single We Do slide."""
        self.add_we_do_slides({
            "examples": [{"question": question, "steps": steps}]
        })
        return self
    
    def add_you_do_slides(self, you_do_data: list):
        """Add You Do slides - FIXED LAYOUT for differentiated tasks."""
        
        # ===== SLIDE 1: BRONZE + SILVER =====
        slide1 = self._add_slide("YOU_DO")
        tf1 = self._create_text_box(slide1, 0.5, 1.3, 12.0, 5.7)
        
        content1 = [("✍️ YOU DO: Independent Practice", True, self.PURPLE), ("\n", False, None)]
        
        for task in you_do_data[:2]:  # Bronze and Silver
            level = task.get('name', 'Task')
            if 'Bronze' in level:
                content1.append(("\n\n🥉 " + level, True, self.ORANGE))
            else:
                content1.append(("\n\n🥈 " + level, True, self.BLUE))
            
            q_text = task.get('question', '')[:120]
            content1.append((f"\n{q_text}", False, None))
            
            # Add answer
            ans = str(task.get('answer', ''))[:60]
            content1.append(("\n✓ Answer: ", False, None))
            content1.append((ans, True, self.GREEN))
        
        self._add_text_content(tf1, content1, font_size=16)
        
        # ===== SLIDE 2: GOLD + EXTENSION =====
        if len(you_do_data) >= 3:
            slide2 = self._add_slide("YOU_DO")
            tf2 = self._create_text_box(slide2, 0.5, 1.3, 12.0, 5.7)
            
            content2 = [("✍️ YOU DO: Challenge Tasks", True, self.PURPLE), ("\n", False, None)]
            
            for task in you_do_data[2:4]:
                level = task.get('name', 'Extension')
                if 'Gold' in level:
                    content2.append(("\n\n🥇 " + level, True, self.RED))
                else:
                    content2.append(("\n\n🚀 " + level, True, self.RED))
                
                q_text = task.get('question', '')[:120]
                content2.append((f"\n{q_text}", False, None))
                
                ans = str(task.get('answer', ''))[:60]
                content2.append(("\n✓ Answer: ", False, None))
                content2.append((ans, True, self.GREEN))
            
            self._add_text_content(tf2, content2, font_size=16)
        
        return self
    
    def add_you_do(self, tasks: list):
        """Legacy You Do method."""
        self.add_you_do_slides(tasks)
        return self
    
    def add_affirmative_checking(self, checkpoint: str = None, action: str = None,
                                  question: str = None, mark_scheme: list = None):
        """Add Affirmative Checking slide - FIXED LAYOUT."""
        slide = self._add_slide("CHECK")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [("✔️ AFFIRMATIVE CHECKING", True, self.RED), ("\n", False, None)]
        
        if checkpoint:
            ckpt = checkpoint[:150]
            content.append(("\n🔍 Check Point: ", True, None))
            content.append((ckpt, False, None))
        
        if action:
            act = action[:150]
            content.append(("\n\n⚡ Action: ", True, None))
            content.append((act, False, None))
        
        if question:
            q = question[:200]
            content.append((f"\n\n📝 {q}", False, None))
        
        if mark_scheme:
            content.append(("\n\n📋 Mark Scheme:", True, None))
            for i, mark in enumerate(mark_scheme[:5], 1):
                if isinstance(mark, dict):
                    content.append((f"\n{i}. {mark.get('step', '')[:50]} = ", False, None))
                    content.append((f"{mark.get('answer', '')[:30]} ✓", True, self.GREEN))
                else:
                    content.append((f"\n{i}. {str(mark)[:80]}", False, None))
        
        self._add_text_content(tf, content, font_size=18)
        return self
    
    def add_exit_ticket(self, question: str, options: list, answer: str):
        """Add Exit Ticket slide - Uses EXIT template layout."""
        slide = self._add_slide("EXIT")
        
        # Find relevant image
        image_path = self._find_image_for_text(question)
        text_width = 8.5 if image_path else 12.0
        
        tf = self._create_text_box(slide, 0.5, 1.3, text_width, 5.7)
        
        q_text = question[:200]
        content = [
            (f"{q_text}", False, None),
            ("\n", False, None)
        ]
        
        for opt in options[:4]:  # Max 4 options
            opt_text = str(opt)[:80]
            content.append((f"\n   {opt_text}", False, None))
        
        ans_text = str(answer)[:60]
        content.append(("\n\n✓ Correct Answer: ", True, None))
        content.append((ans_text, True, self.GREEN))
        
        self._add_text_content(tf, content, font_size=18)
        
        if image_path:
            self._add_image_to_slide(slide, image_path, left=9.5, top=1.5, width=3.3)
        
        return self
    
    # =========================================================================
    # TLAG TECHNIQUE SLIDES
    # =========================================================================
    
    def add_cold_call(self, question: str, answer: str = None):
        """Add a Cold Call slide using the COLD_CALL template."""
        slide = self._add_slide("COLD_CALL")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [(f"📝 {question[:180]}", False, None)]
        
        if answer:
            content.append(("\n\n✓ Expected Answer: ", True, None))
            content.append((str(answer)[:100], True, self.GREEN))
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_show_me(self, instruction: str, what_to_show: str = None):
        """Add a Show Me slide using the SHOW_ME template."""
        slide = self._add_slide("SHOW_ME")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [(f"📝 {instruction[:180]}", False, None)]
        
        if what_to_show:
            content.append((f"\n\n👆 Show: {what_to_show[:100]}", True, self.BLUE))
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_everybody_writes(self, prompt: str, time_mins: int = 2):
        """Add an Everybody Writes slide using the EVERYBODY_WRITES template."""
        slide = self._add_slide("EVERYBODY_WRITES")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [
            (f"📝 {prompt[:200]}", False, None),
            (f"\n\n⏱️ Time: {time_mins} minutes", True, self.ORANGE)
        ]
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_show_call(self, question: str, expected_response: str = None):
        """Add a Show Call slide using the SHOW_CALL template."""
        slide = self._add_slide("SHOW_CALL")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [(f"📝 {question[:180]}", False, None)]
        
        if expected_response:
            content.append(("\n\n✓ Look for: ", True, None))
            content.append((str(expected_response)[:100], True, self.GREEN))
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_turn_and_talk(self, discussion_prompt: str, time_seconds: int = 30):
        """Add a Turn and Talk slide using the TURN_TALK template."""
        slide = self._add_slide("TURN_TALK")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [
            (f"📝 {discussion_prompt[:200]}", False, None),
            (f"\n\n⏱️ Time: {time_seconds} seconds", True, self.ORANGE)
        ]
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_stretch_it(self, follow_up_question: str, hint: str = None):
        """Add a Stretch It slide using the STRETCH_IT template."""
        slide = self._add_slide("STRETCH_IT")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [(f"🤔 {follow_up_question[:180]}", False, None)]
        
        if hint:
            content.append((f"\n\n💡 Hint: {hint[:100]}", False, self.GRAY))
        
        self._add_text_content(tf, content, font_size=20)
        return self
    
    def add_right_is_right(self, question: str, partial_answer: str, correct_answer: str):
        """Add a Right is Right slide using the RIGHT_IS_RIGHT template."""
        slide = self._add_slide("RIGHT_IS_RIGHT")
        tf = self._create_text_box(slide, 0.5, 1.3, 12.0, 5.7)
        
        content = [
            (f"📝 {question[:150]}", False, None),
            ("\n\n❌ Partial answer: ", True, None),
            (f"{partial_answer[:80]}", False, self.ORANGE),
            ("\n\n✓ Complete answer: ", True, None),
            (f"{correct_answer[:80]}", True, self.GREEN)
        ]
        
        self._add_text_content(tf, content, font_size=18)
        return self
    
    def save(self, output_path: str):
        """Save the presentation to file."""
        self.prs.save(output_path)
        return output_path
    
    def save_to_bytes(self) -> BytesIO:
        """Save the presentation to BytesIO for direct download."""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer
